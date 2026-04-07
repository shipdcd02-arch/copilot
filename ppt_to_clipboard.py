"""
Script 2: ppt_to_clipboard.py
clipboard_to_ppt.exe 가 생성한 임시 PPT 파일 경로를 인자로 받아,
텍스트를 읽어 클립보드에 넣고, 파일을 자동 삭제합니다.

필요 패키지: pip install pyperclip python-pptx
exe 빌드:   pyinstaller --onefile --noconsole ppt_to_clipboard.py
"""

import sys
import os

import pyperclip
from pptx import Presentation


def main():
    if len(sys.argv) < 2:
        print("사용법: ppt_to_clipboard.exe <pptx_file_path>")
        sys.exit(1)

    pptx_path = sys.argv[1]

    if not os.path.exists(pptx_path):
        print(f"파일을 찾을 수 없습니다: {pptx_path}")
        sys.exit(1)

    try:
        # 1. PPT에서 텍스트 읽기
        prs = Presentation(pptx_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)

        full_text = "\n".join(parts)

        # 2. 클립보드에 넣기
        pyperclip.copy(full_text)

    finally:
        # 3. 임시 파일 삭제 (성공/실패 무관하게 항상 삭제)
        try:
            os.remove(pptx_path)
        except Exception:
            pass


if __name__ == "__main__":
    main()
