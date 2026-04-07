"""
Script 1: clipboard_to_ppt.py
클립보드의 텍스트를 임시 PPT 파일로 저장하고,
ppt_to_clipboard.exe를 호출하여 다시 클립보드로 읽어오게 합니다.

필요 패키지: pip install pyperclip python-pptx
exe 빌드:   pyinstaller --onefile --noconsole clipboard_to_ppt.py
"""

import sys
import os
import tempfile
import subprocess

import pyperclip
from pptx import Presentation
from pptx.util import Inches, Pt


def main():
    # 1. 클립보드에서 텍스트 읽기
    text = pyperclip.paste()
    if not text.strip():
        print("클립보드가 비어있습니다.")
        sys.exit(1)

    # 2. 임시 PPT 파일 생성 (한글 경로 방지 → ASCII 경로 사용)
    temp_dir = tempfile.gettempdir()
    try:
        temp_dir.encode("ascii")
    except UnicodeEncodeError:
        # 기본 temp 폴더에 한글이 있으면 C:\Temp 사용
        temp_dir = "C:\\Temp"
        os.makedirs(temp_dir, exist_ok=True)

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False, dir=temp_dir)
    tmp_path = tmp.name
    tmp.close()

    # 3. PPT에 텍스트 저장
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = text

    prs.save(tmp_path)

    # 4. ppt_to_clipboard.exe 경로 결정
    #    - exe 실행 시: 자신과 같은 폴더에서 찾음
    #    - .py 실행 시: 이 스크립트와 같은 폴더에서 찾음
    if getattr(sys, "frozen", False):
        base_dir = os.path.dirname(sys.executable)
    else:
        base_dir = os.path.dirname(os.path.abspath(__file__))

    exe_path = os.path.join(base_dir, "ppt_to_clipboard.exe")

    if not os.path.exists(exe_path):
        print(f"ppt_to_clipboard.exe 를 찾을 수 없습니다: {exe_path}")
        # 임시파일 정리
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        sys.exit(1)

    # 5. ppt_to_clipboard.exe에 임시 파일 경로 전달 (비동기 실행)
    #    한글 경로 + --noconsole exe 호환을 위해 shell=True + 큰따옴표 처리
    cmd = f'"{exe_path}" "{tmp_path}"'
    subprocess.Popen(
        cmd,
        shell=True,
        stdin=subprocess.DEVNULL,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


if __name__ == "__main__":
    main()
